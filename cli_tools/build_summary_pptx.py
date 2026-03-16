#!/usr/bin/env python3
"""
Build a summary PowerPoint from a Yad2 pipeline output directory.

Each slide = one exported listing:
- Title: city + address + price
- Subtitle: original URL
- Right: debug screenshot PNG (listing page)
- Left upper: listing images + description text box
- Left bottom: floor, parking, transportation, property status

Usage:
  python scripts/build_summary_pptx.py --output-dir output
  python scripts/build_summary_pptx.py --output-dir output --out-pptx summary.pptx
"""

import argparse
import re
import sys
from pathlib import Path
from typing import List, Optional

import pandas as pd
from pptx import Presentation
from pptx.dml.color import RgbColor
from pptx.util import Inches, Pt

# Slide dimensions (default 13.333 x 7.5 in)
SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)
MARGIN = Inches(0.4)
TITLE_TOP = Inches(0.2)
TITLE_HEIGHT = Inches(0.6)
SUBTITLE_TOP = Inches(0.75)
SUBTITLE_HEIGHT = Inches(0.35)
LEFT_WIDTH = Inches(7.8)
RIGHT_LEFT = Inches(8.2)
RIGHT_WIDTH = Inches(4.9)
CONTENT_TOP = Inches(1.15)
CONTENT_HEIGHT = Inches(6.0)
LEFT_IMAGES_HEIGHT = Inches(3.2)
LEFT_DETAILS_TOP = Inches(4.4)


def _find_debug_png(debug_dir: Path, listing_id: str) -> Optional[Path]:
    """Return path to the most recent debug PNG for this listing (exported screenshot)."""
    pattern = f"{re.escape(str(listing_id))}_*.png"
    matches = sorted(debug_dir.glob(pattern), key=lambda p: p.name, reverse=True)
    return matches[0] if matches else None


def _find_listing_images(images_dir: Path, listing_id: str) -> List[Path]:
    """Return list of image paths for this listing, in order."""
    folder = images_dir / str(listing_id)
    if not folder.is_dir():
        return []
    return sorted(folder.glob("*.*"), key=lambda p: p.name)


def _safe_str(val, default: str = "—") -> str:
    if val is None or (isinstance(val, float) and (val != val)):
        return default
    try:
        if hasattr(pd, "isna") and pd.isna(val):
            return default
    except Exception:
        pass
    return str(val).strip() or default


def _format_price(price) -> str:
    if price is None or (isinstance(price, float) and (price != price)):
        return "—"
    try:
        return f"{int(float(price)):,} ₪".replace(",", " ")
    except (ValueError, TypeError):
        return str(price)


def build_pptx(output_dir: Path, out_pptx: Path) -> None:
    output_dir = Path(output_dir).resolve()
    csv_path = output_dir / "listings_full.csv"
    debug_dir = output_dir / "debug"
    images_dir = output_dir / "images"

    if not csv_path.exists():
        print(f"CSV not found: {csv_path}", file=sys.stderr)
        sys.exit(1)

    df = pd.read_csv(csv_path, encoding="utf-8")
    if df.empty:
        print("No listings in CSV.", file=sys.stderr)
        sys.exit(0)

    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    for _, row in df.iterrows():
        listing_id = _safe_str(row.get("yad2_listing_id"), "")
        if not listing_id:
            continue

        city = _safe_str(row.get("city"))
        address = _safe_str(row.get("full_address_best") or row.get("street_name") or row.get("street_address_raw"))
        price = _format_price(row.get("price_ils"))
        url = _safe_str(row.get("original_listing_url"))

        title_text = f"{city}, {address} — {price}"
        if title_text.startswith("—"):
            title_text = f"Listing {listing_id} — {price}"

        # Blank slide
        blank = prs.slide_layouts[6]  # blank
        slide = prs.slides.add_slide(blank)

        # Title
        title_box = slide.shapes.add_textbox(MARGIN, TITLE_TOP, SLIDE_WIDTH - 2 * MARGIN, TITLE_HEIGHT)
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title_text[:200]
        p.font.size = Pt(18)
        p.font.bold = True

        # Subtitle (URL)
        sub_box = slide.shapes.add_textbox(MARGIN, SUBTITLE_TOP, SLIDE_WIDTH - 2 * MARGIN, SUBTITLE_HEIGHT)
        tf = sub_box.text_frame
        p = tf.paragraphs[0]
        p.text = url[:250]
        p.font.size = Pt(9)
        p.font.color.rgb = RgbColor(0x44, 0x44, 0x44)

        # Right: debug PNG (screenshot of listing page, saved only for exported listings)
        png_path = _find_debug_png(debug_dir, listing_id) if debug_dir.exists() else None
        if png_path and png_path.exists():
            try:
                slide.shapes.add_picture(
                    str(png_path),
                    RIGHT_LEFT,
                    CONTENT_TOP,
                    width=RIGHT_WIDTH,
                    height=CONTENT_HEIGHT,
                )
            except Exception:
                pass

        # Left upper: description text (and optionally first image)
        desc_text = _safe_str(row.get("property_summary") or row.get("description_clean"), "")
        if desc_text:
            desc_text = (desc_text[:600] + "…") if len(desc_text) > 600 else desc_text
        img_paths = _find_listing_images(images_dir, listing_id)
        img_height = Inches(2.0)
        img_left = MARGIN
        for i, ip in enumerate(img_paths[:3]):  # up to 3 images
            try:
                slide.shapes.add_picture(
                    str(ip),
                    img_left,
                    CONTENT_TOP + Inches(0.1),
                    width=Inches(2.2),
                    height=img_height,
                )
                img_left += Inches(2.3)
            except Exception:
                pass
        desc_box = slide.shapes.add_textbox(
            MARGIN,
            CONTENT_TOP + img_height + Inches(0.15),
            LEFT_WIDTH - Inches(0.2),
            LEFT_IMAGES_HEIGHT - img_height - Inches(0.1),
        )
        tf = desc_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = desc_text or "No description."
        p.font.size = Pt(9)

        # Left bottom: floor, parking, transportation, property status
        floor_cur = row.get("floor_current")
        floor_tot = row.get("floor_total")
        floor_str = "—"
        if floor_cur is not None and floor_tot is not None:
            floor_str = f"Floor {int(floor_cur)} of {int(floor_tot)}"
        elif floor_cur is not None:
            floor_str = f"Floor {int(floor_cur)}"

        parking = row.get("parking_count")
        parking_str = f"Parking: {int(parking)}" if parking is not None else "Parking: —"

        ta_min = row.get("drive_to_tel_aviv_savidor_duration_min")
        bs_min = row.get("drive_to_beer_sheva_center_duration_min")
        transport_parts = []
        if ta_min is not None:
            transport_parts.append(f"→ TLV Savidor: {float(ta_min):.0f} min")
        if bs_min is not None:
            transport_parts.append(f"→ Beer Sheva: {float(bs_min):.0f} min")
        commute = _safe_str(row.get("commute_assessment"), "")
        if commute:
            transport_parts.append(commute[:120])
        transport_str = "\n".join(transport_parts) if transport_parts else "Transport: —"

        prop_status = _safe_str(row.get("property_condition_label"), "")

        details_lines = [
            f"Floor: {floor_str}",
            parking_str,
            "Transportation:",
            transport_str,
            f"Property: {prop_status}" if prop_status else "",
        ]
        details_text = "\n".join(l for l in details_lines if l).strip()

        details_box = slide.shapes.add_textbox(
            MARGIN,
            LEFT_DETAILS_TOP,
            LEFT_WIDTH,
            Inches(2.8),
        )
        tf = details_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = details_text
        p.font.size = Pt(10)

    out_pptx = Path(out_pptx).resolve()
    out_pptx.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out_pptx))
    print(f"Saved: {out_pptx} ({len(df)} slides)")


def main() -> None:
    parser = argparse.ArgumentParser(
        description="Build a summary PowerPoint from a Yad2 pipeline output directory."
    )
    parser.add_argument(
        "--output-dir",
        type=str,
        default="output",
        help="Pipeline output directory (contains listings_full.csv, debug/, images/)",
    )
    parser.add_argument(
        "--out-pptx",
        type=str,
        default=None,
        help="Output PowerPoint path (default: <output-dir>/summary_listings.pptx)",
    )
    args = parser.parse_args()

    output_dir = Path(args.output_dir)
    out_pptx = args.out_pptx or str(output_dir / "summary_listings.pptx")
    build_pptx(output_dir, out_pptx)


if __name__ == "__main__":
    main()
